from typing import List
from fastapi import FastAPI, UploadFile, File, HTTPException
from PyPDF2 import PdfReader
import docx
import os
from bs4 import BeautifulSoup
import pypandoc
from pptx import Presentation
from nltk.tokenize import sent_tokenize
import re
from tempfile import TemporaryFile
from pydantic import BaseModel


app = FastAPI(
    title="Text Extractor",
    description="Extract text from a document. "
    "Supported file types include PDF, DOCX, PPTX, TXT, HTML, and RTF.",
    version="0.0.1",
)


class ExtractedTextResponse(BaseModel):
    text: str


class ChunkResponse(BaseModel):
    chunks: List[str]


@app.post(
    "/extract-text/",
    response_model=ExtractedTextResponse,
    summary="Extract text from a file",
    description="Upload a file and extract its text content.",
)
async def extract_text(file: UploadFile = File(...)):
    text = await extract_text_from_file(file)
    return {"text": text}


@app.post(
    "/extract-chunks/",
    response_model=ChunkResponse,
    summary="Extract text chunks from a file",
    description="Upload a file, extract its text content, "
    "and divide it into chunks based on the specified number of "
    "sentences per chunk and overlap.",
)
async def extract_chunk(
    file: UploadFile = File(...), sentences_per_chunk: int = 5, overlap: int = 0
):
    text = await extract_text_from_file(file)

    if not text.strip():
        raise HTTPException(
            status_code=400,
            detail="The uploaded file is empty or not readable.",
        )

    try:
        chunks = chunk_text_by_sentences(text, sentences_per_chunk, overlap)
        return {"chunks": chunks}
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))


def extract_text_from_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text


def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text


def extract_text_from_pptx(file):
    with TemporaryFile() as tmp_file:
        tmp_file.write(file.read())
        tmp_file.seek(0)
        prs = Presentation(tmp_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text


def extract_text_from_txt(file):
    return file.read().decode("utf-8")


def extract_text_from_html(file):
    soup = BeautifulSoup(file.read(), "html.parser")
    return soup.get_text()


def extract_text_from_rtf(file):
    output = pypandoc.convert_text(file.read().decode("utf-8"), "plain", format="rtf")
    return output


async def extract_text_from_file(file: UploadFile) -> str:
    file_extension = os.path.splitext(file.filename)[-1].lower()

    if file_extension == ".pdf":
        return extract_text_from_pdf(file.file)
    elif file_extension == ".docx":
        return extract_text_from_docx(file.file)
    elif file_extension == ".pptx":
        return extract_text_from_pptx(file.file)
    elif file_extension == ".txt":
        return extract_text_from_txt(await file.read())
    elif file_extension == ".html":
        return extract_text_from_html(await file.read())
    elif file_extension == ".rtf":
        return extract_text_from_rtf(await file.read())
    else:
        return (await file.read()).decode("utf-8")


def chunk_text_by_sentences(
    source_text: str, sentences_per_chunk: int, overlap: int
) -> list:
    if sentences_per_chunk < 2:
        raise ValueError("Sentences per chunk must be 2 or more.")
    if overlap < 0 or overlap >= sentences_per_chunk:
        raise ValueError(
            "Overlap must be 0 or more and less than the number of sentences."
        )

    cleaned_text = re.sub(r"\s+", " ", source_text.replace("\n", " ").strip())

    sentences = sent_tokenize(cleaned_text)
    if not sentences:
        print("Nothing to chunk")
        return []

    chunks = []
    i = 0

    while i < len(sentences):
        end = min(i + sentences_per_chunk, len(sentences))
        chunk = " ".join(sentences[i:end])

        if overlap > 0 and i > 0:
            overlap_start = max(0, i - overlap)
            overlap_chunk = " ".join(sentences[overlap_start:i])
            chunk = overlap_chunk + " " + chunk

        chunks.append(chunk.strip())
        i += sentences_per_chunk - overlap

    return chunks
